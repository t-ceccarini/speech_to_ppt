import openai
from openai import OpenAI
from pptx import Presentation
import os

def generate_presentation_from_transcription(api_key, transcription, output_file):
    """
    Genera una presentazione PowerPoint basata su una trascrizione audio.

    :param transcription: Testo della trascrizione audio
    :param output_file: Nome del file PowerPoint di output
    """
    # Imposta la chiave API di OpenAI
    openai.api_key = api_key

    # Prompt dinamico: chiediamo a GPT di strutturare una presentazione sulla base della trascrizione
    prompt = f"""
    Crea il contenuto di una presentazione PowerPoint basata sulla seguente trascrizione:
    ---
    {transcription}
    ---
    Identifica i temi principali e organizza le informazioni in una serie di diapositive.
    Ogni diapositiva deve avere:
    - Un titolo (chiaro e breve)
    - Contenuti in formato elenco puntato
    Rendi la struttura della presentazione coerente con il contenuto della trascrizione.
    """

    # Richiesta all'API OpenAI
    response = openai.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "Sei un esperto nel creare presentazioni PowerPoint professionali."},
            {"role": "user", "content": prompt}
        ]
    )

    # Estrai il contenuto generato
    generated_content = response.choices[0].message.content

    # Divide il testo in sezioni per ogni diapositiva
    slides_content = generated_content.split("\n\n")

    # Crea la presentazione PowerPoint
    prs = Presentation()

    for slide_text in slides_content:
        # Divide titolo e contenuto
        lines = slide_text.split("\n")
        if len(lines) > 1:
            title = lines[0].strip()  # Il titolo è la prima riga
            content = "\n".join(lines[1:]).strip()  # Il resto è il contenuto
        else:
            title = "Slide senza titolo"
            content = slide_text.strip()

        # Aggiungi una diapositiva con titolo e contenuto
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout con titolo e contenuto
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    # Salva la presentazione
    prs.save(output_file)
    print(f"Presentazione creata con successo: {output_file}")

def generate_presentation_with_template(template_path, slides_content, output_path):
    """
    Genera una presentazione utilizzando un template esistente.

    :param template_path: Path del file template .pptx
    :param slides_content: Lista di dizionari con contenuti delle slide.
                           Es: [{"type": "title", "title": "...", "content": "..."}, ...]
    :param output_path: Path completo del file PowerPoint generato (.pptx)
    """
    # Controlla se il template esiste
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Il template specificato non esiste: {template_path}")

    # Carica il file template
    prs = Presentation(template_path)

    # Aggiungi le slide una per una
    for slide_data in slides_content:
        # Determina il layout basato sul tipo della slide
        if slide_data["type"] == "title":
            slide_layout = prs.slide_layouts[0]  # Layout per "Slide Titolo"
        elif slide_data["type"] == "content":
            slide_layout = prs.slide_layouts[1]  # Layout per "Titolo + Contenuto"
        elif slide_data["type"] == "graph":
            slide_layout = prs.slide_layouts[5]  # Layout per "Grafico o immagine"
        else:
            raise ValueError(f"Tipo di slide sconosciuto: {slide_data['type']}")

        # Crea una nuova slide con il layout scelto
        slide = prs.slides.add_slide(slide_layout)
        
        print(slide)
        print(slide_data)
        print(slide_data["title"])

        # Aggiungi titolo alla slide
        if "title" in slide_data:
            slide.shapes.title.text = slide_data["title"]

        # Aggiungi contenuto
        if "content" in slide_data and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data["content"]

    # Salva la presentazione
    prs.save(output_path)
    print(f"Presentazione generata con successo: {output_path}")